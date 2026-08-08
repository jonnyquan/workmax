#!/usr/bin/env python3
"""validate-pptx.py — structural + content check for generated PPTX files.

Invoked by the Go agent runtime after a ppt-mode turn completes; sees
the .pptx the SDK wrote and emits a JSON status the runtime attaches
to the user-facing done event.

Usage:
    python3 validate-pptx.py <path/to/file.pptx>

Output (stdout):
    {"status": "ok"|"warning"|"error", "issues": ["..."], "slide_count": N}

Exit code: always 0. Status lives in the JSON body — the Go runner
parses stdout and never branches on the exit code. (subprocess.run
with check=False keeps Python tracebacks out of stderr noise.)

Graceful degradation:
    The structural checks (file is a zip, ppt/presentation.xml present,
    >= 1 slide) need only the stdlib. Content checks (placeholder
    detection, empty slides) need python-pptx. When python-pptx is not
    installed we still emit a useful status — the structural failure
    modes are the most common ones in production anyway.
"""

import json
import re
import sys
import zipfile

# Markers we treat as "the model gave up halfway and wrote a placeholder
# instead of real content". These cover the most common LLM lazy
# patterns plus the legacy lorem-ipsum fingerprint. Add to the list
# as new failure modes surface in production logs.
_PLACEHOLDER_PATTERNS = [
    r"\[TODO",
    r"\[PLACEHOLDER",
    r"\[INSERT[ _-]?HERE",
    r"\[FILL[ _-]?IN",
    r"\bTBD\b",
    r"lorem ipsum",
]


def _emit(status, issues, **extra):
    """Single-line JSON output — keeps the Go side's parser dumb."""
    payload = {"status": status, "issues": issues}
    payload.update(extra)
    sys.stdout.write(json.dumps(payload, ensure_ascii=False))


def main():
    if len(sys.argv) < 2:
        _emit("error", ["missing pptx path argument"])
        return

    pptx_path = sys.argv[1]

    # Step 1: file opens as a zip. Catches everything from "doesn't
    # exist" through "0-byte file" through "corrupted output stream".
    try:
        zf = zipfile.ZipFile(pptx_path)
    except FileNotFoundError:
        _emit("error", ["file not found: " + pptx_path])
        return
    except zipfile.BadZipFile as exc:
        _emit("error", [f"file is not a valid zip: {exc}"])
        return
    except Exception as exc:  # IO error / permission / etc.
        _emit("error", [f"unable to open file: {exc}"])
        return

    names = zf.namelist()

    # Step 2: looks like a real PPTX. The presentation.xml entry is
    # the openxml signature for the format — without it we have a
    # zip that happens to be named .pptx (PowerPoint refuses to open).
    if "ppt/presentation.xml" not in names:
        _emit("error", ["not a valid PPTX (ppt/presentation.xml missing)"])
        return

    slide_files = sorted(
        n for n in names if re.match(r"ppt/slides/slide\d+\.xml$", n)
    )
    if not slide_files:
        _emit("error", ["zero slides — the SDK wrote an empty presentation"])
        return

    # Step 3: try the deeper content checks. python-pptx is a soft
    # dependency — when missing we still give the runtime a useful
    # status (the structural checks above succeeded, so the file is
    # at least openable; we just couldn't introspect contents).
    try:
        from pptx import Presentation
    except ImportError:
        _emit(
            "ok",
            [],
            slide_count=len(slide_files),
            checks_skipped="python-pptx not installed; structural-only validation",
        )
        return

    try:
        prs = Presentation(pptx_path)
    except Exception as exc:
        _emit("error", [f"python-pptx parse failed: {exc}"])
        return

    placeholder_re = re.compile("|".join(_PLACEHOLDER_PATTERNS), re.IGNORECASE)
    issues = []

    for idx, slide in enumerate(prs.slides, start=1):
        # Collect every text run on the slide so the placeholder
        # regex sees the entire slide content as one string. Iterating
        # at the run level instead of shape level catches placeholder
        # text that lives inside text frames with formatting.
        texts = []
        has_text_content = False
        has_picture = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text = shape.text_frame.text
                if slide_text and slide_text.strip():
                    texts.append(slide_text)
                    has_text_content = True
            # MSO_SHAPE_TYPE.PICTURE = 13 (avoid importing the enum
            # to keep the import surface minimal — the integer is
            # part of the openxml spec, not a python-pptx invention)
            try:
                if shape.shape_type == 13:
                    has_picture = True
            except Exception:
                # Some shapes raise when we read shape_type; ignore.
                pass

        joined = "\n".join(texts)
        if joined and placeholder_re.search(joined):
            issues.append(f"slide {idx} contains placeholder marker(s)")

        if not has_text_content and not has_picture:
            issues.append(f"slide {idx} has no visible content")

    status = "warning" if issues else "ok"
    _emit(status, issues, slide_count=len(prs.slides))


if __name__ == "__main__":
    main()
